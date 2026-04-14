"""
ADAM DNA Tool - Document Ingestion & Processing
Handles file uploads (DOCX, PPTX, PDF, CSV, JSON, XLSX) and URL fetching.
Extracts text, identifies relevant DNA sections, and maps content to the questionnaire structure.
"""

import os
import io
import json
import re
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime
import structlog
import aiohttp
import aiofiles

logger = structlog.get_logger()


class DocumentProcessor:
    """Processes uploaded documents and extracts text + structural information."""

    SUPPORTED_TYPES = {
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".pdf": "application/pdf",
        ".csv": "text/csv",
        ".json": "application/json",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".yaml": "application/x-yaml",
        ".yml": "application/x-yaml",
    }

    # Keywords that indicate relevance to specific DNA sections
    SECTION_KEYWORDS = {
        "1": ["mission", "vision", "principles", "values", "boundaries", "governance",
              "directors", "board", "constitution", "founding", "charter", "incorporated"],
        "2": ["culture", "norms", "behavior", "failure", "brand voice", "internal",
              "external", "posture", "identity", "decision making", "consensus"],
        "3": ["objectives", "goals", "mandates", "okr", "kpi", "targets",
              "strategic priorities", "fiscal year", "aspiration"],
        "4": ["rules", "regulations", "compliance", "expectations", "zero tolerance",
              "exceptions", "policy", "constraints", "prohibition"],
        "5": ["budget", "revenue", "financial", "rights", "licensing", "ip",
              "customer segment", "regulatory", "jurisdiction", "strategy drift",
              "reputation", "supplier", "partner"],
        "6": ["boss", "scoring", "risk", "threshold", "escalation", "exception",
              "sovereignty score", "weighting", "dimensions"],
        "7": ["intent", "conflict", "arbitration", "risk tolerance", "urgency",
              "approval", "delegation", "authorization"],
        "8": ["agent", "domain governor", "work group", "digital twin",
              "architecture", "orchestration", "mesh", "ai centric"],
        "9": ["flight recorder", "evidence", "audit", "forensic", "retention",
              "tamper", "hash chain", "compliance reporting"],
        "10": ["products", "services", "revenue model", "customer lifecycle",
               "competitive", "supply chain", "pricing", "market"],
        "11": ["seasonal", "temporal", "regional", "timezone", "variance",
               "refresh cadence", "jurisdictional override"],
        "12": ["cloud", "azure", "aws", "gcp", "kubernetes", "infrastructure",
               "compute", "storage", "networking", "sovereignty", "data residency",
               "on premises", "hybrid", "multi-cloud"],
        "13": ["resilience", "idempotency", "security", "disaster", "failover",
               "encryption", "zero trust", "threat model", "backup", "recovery"],
    }

    def __init__(self, upload_dir: str):
        self.upload_dir = upload_dir
        os.makedirs(upload_dir, exist_ok=True)

    async def process_file(self, filename: str, content: bytes) -> Dict[str, Any]:
        """Process an uploaded file and return extracted information."""
        ext = os.path.splitext(filename)[1].lower()

        if ext not in self.SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: {ext}. Supported: {', '.join(self.SUPPORTED_TYPES.keys())}")

        # Save file
        filepath = os.path.join(self.upload_dir, filename)
        async with aiofiles.open(filepath, "wb") as f:
            await f.write(content)

        # Extract text based on file type
        extracted_text = ""
        extracted_sections = {}

        if ext == ".docx":
            extracted_text, extracted_sections = self._process_docx(filepath)
        elif ext == ".pptx":
            extracted_text, extracted_sections = self._process_pptx(filepath)
        elif ext == ".pdf":
            extracted_text, extracted_sections = self._process_pdf(filepath)
        elif ext == ".csv":
            extracted_text = self._process_csv(filepath)
        elif ext == ".json":
            extracted_text = self._process_json(filepath)
        elif ext == ".xlsx":
            extracted_text = self._process_xlsx(filepath)
        elif ext in (".txt", ".md"):
            extracted_text = self._process_text(filepath)
        elif ext in (".yaml", ".yml"):
            extracted_text = self._process_yaml(filepath)

        # Calculate relevance mapping
        relevance = self._calculate_relevance(extracted_text)

        return {
            "filename": filename,
            "filepath": filepath,
            "mime_type": self.SUPPORTED_TYPES.get(ext, "application/octet-stream"),
            "size_bytes": len(content),
            "extracted_text": extracted_text[:50000],  # Cap at 50K chars
            "extracted_sections": extracted_sections,
            "relevance_mapping": relevance,
        }

    async def fetch_url(self, url: str) -> Dict[str, Any]:
        """Fetch content from a URL and process it."""
        logger.info("Fetching URL", url=url)

        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=aiohttp.ClientTimeout(total=30)) as response:
                if response.status != 200:
                    raise ValueError(f"Failed to fetch URL: HTTP {response.status}")

                content_type = response.headers.get("Content-Type", "").lower()
                content = await response.read()

                # Determine file extension from content type
                ext = ".txt"
                if "json" in content_type:
                    ext = ".json"
                elif "csv" in content_type:
                    ext = ".csv"
                elif "yaml" in content_type or "yml" in content_type:
                    ext = ".yaml"
                elif "pdf" in content_type:
                    ext = ".pdf"
                elif "html" in content_type:
                    ext = ".html"

                filename = f"url_fetch_{datetime.now().strftime('%Y%m%d_%H%M%S')}{ext}"

                if ext == ".html":
                    # Extract text from HTML
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(content, "html.parser")
                    for script in soup(["script", "style"]):
                        script.decompose()
                    text = soup.get_text(separator="\n", strip=True)
                    relevance = self._calculate_relevance(text)

                    return {
                        "filename": filename,
                        "source_url": url,
                        "mime_type": content_type,
                        "size_bytes": len(content),
                        "extracted_text": text[:50000],
                        "extracted_sections": {},
                        "relevance_mapping": relevance,
                    }

                return await self.process_file(filename, content)

    def _process_docx(self, filepath: str) -> Tuple[str, Dict]:
        """Extract text and structure from a DOCX file."""
        from docx import Document
        doc = Document(filepath)

        full_text = []
        sections = {}
        current_heading = "Introduction"

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            if para.style.name.startswith("Heading"):
                current_heading = text
                sections[current_heading] = []

            full_text.append(text)
            if current_heading in sections:
                sections[current_heading].append(text)

        # Also extract table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    full_text.append(row_text)

        return "\n".join(full_text), sections

    def _process_pptx(self, filepath: str) -> Tuple[str, Dict]:
        """Extract text from a PPTX presentation."""
        from pptx import Presentation
        prs = Presentation(filepath)

        full_text = []
        sections = {}

        for i, slide in enumerate(prs.slides, 1):
            slide_key = f"Slide {i}"
            slide_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if slide_texts:
                sections[slide_key] = slide_texts
                full_text.extend(slide_texts)

            # Extract from tables
            if shape.has_table if hasattr(shape, "has_table") else False:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        full_text.append(row_text)

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    full_text.append(f"[Speaker Notes] {notes}")

        return "\n".join(full_text), sections

    def _process_pdf(self, filepath: str) -> Tuple[str, Dict]:
        """Extract text from a PDF file."""
        from pypdf import PdfReader
        reader = PdfReader(filepath)

        full_text = []
        sections = {}

        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                sections[f"Page {i}"] = [text.strip()]
                full_text.append(text.strip())

        return "\n".join(full_text), sections

    def _process_csv(self, filepath: str) -> str:
        """Extract text from a CSV file."""
        import pandas as pd
        df = pd.read_csv(filepath, nrows=1000)
        return f"Columns: {', '.join(df.columns)}\n\nData Sample:\n{df.head(50).to_string()}\n\nShape: {df.shape}"

    def _process_xlsx(self, filepath: str) -> str:
        """Extract text from an Excel file."""
        import pandas as pd
        xls = pd.ExcelFile(filepath)
        parts = []
        for sheet in xls.sheet_names[:10]:  # Limit to 10 sheets
            df = pd.read_excel(filepath, sheet_name=sheet, nrows=500)
            parts.append(f"=== Sheet: {sheet} ===\nColumns: {', '.join(str(c) for c in df.columns)}\n{df.head(30).to_string()}\n")
        return "\n".join(parts)

    def _process_json(self, filepath: str) -> str:
        """Extract text from a JSON file."""
        with open(filepath, "r") as f:
            data = json.load(f)
        return json.dumps(data, indent=2)[:50000]

    def _process_text(self, filepath: str) -> str:
        """Read plain text or markdown file."""
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read()[:50000]

    def _process_yaml(self, filepath: str) -> str:
        """Read YAML file."""
        import yaml
        with open(filepath, "r") as f:
            data = yaml.safe_load(f)
        return json.dumps(data, indent=2, default=str)[:50000]

    def _calculate_relevance(self, text: str) -> Dict[str, float]:
        """Calculate relevance scores mapping text to DNA sections."""
        text_lower = text.lower()
        relevance = {}

        for section_num, keywords in self.SECTION_KEYWORDS.items():
            score = 0.0
            for keyword in keywords:
                count = text_lower.count(keyword.lower())
                if count > 0:
                    score += min(count * 0.1, 1.0)  # Cap per-keyword contribution

            # Normalize to 0-1
            max_possible = len(keywords)
            relevance[section_num] = round(min(score / max(max_possible * 0.3, 1), 1.0), 2)

        return relevance

    def get_relevant_sections(self, relevance: Dict[str, float], threshold: float = 0.2) -> List[str]:
        """Get DNA sections that a document is relevant to (above threshold)."""
        return sorted(
            [s for s, score in relevance.items() if score >= threshold],
            key=lambda s: relevance[s],
            reverse=True,
        )
