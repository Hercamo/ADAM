"""ADAM Book corpus loader.

The connector can point at a local copy of the ADAM Book directory and expose
its documents to the AI as read-only context. Supports .md / .txt directly
and .docx / .pdf via lazy, optional dependencies.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

log = logging.getLogger("adam.corpus")

TEXT_EXTS = {".md", ".txt", ".yaml", ".yml", ".json", ".py", ".html"}
DOCX_EXT = ".docx"
PDF_EXT = ".pdf"
PPTX_EXT = ".pptx"

def _walk_corpus(root: Path) -> List[Path]:
    out: List[Path] = []
    if not root.exists():
        return out
    for p in root.rglob("*"):
        if p.is_file() and not p.name.startswith(".") and "__pycache__" not in p.parts:
            out.append(p)
    return out

def list_documents(corpus_dir: Optional[str]) -> List[Dict[str, Any]]:
    if not corpus_dir:
        return []
    root = Path(corpus_dir)
    docs: List[Dict[str, Any]] = []
    for p in _walk_corpus(root):
        ext = p.suffix.lower()
        if ext not in TEXT_EXTS and ext not in {DOCX_EXT, PDF_EXT, PPTX_EXT}:
            continue
        try:
            size = p.stat().st_size
        except OSError:
            size = 0
        docs.append({
            "path": str(p.relative_to(root)),
            "ext": ext,
            "size": size,
            "readable": ext in TEXT_EXTS or ext == DOCX_EXT or ext == PDF_EXT,
        })
    return docs

def read_document(rel_path: str, corpus_dir: Optional[str], max_chars: int = 20000) -> Dict[str, Any]:
    if not corpus_dir:
        return {"path": rel_path, "error": "No corpus_dir configured."}
    root = Path(corpus_dir).resolve()
    target = (root / rel_path).resolve()
    # prevent path traversal
    try:
        target.relative_to(root)
    except ValueError:
        return {"path": rel_path, "error": "Path escapes corpus root."}

    if not target.exists():
        return {"path": rel_path, "error": "Not found."}

    ext = target.suffix.lower()
    try:
        if ext in TEXT_EXTS:
            text = target.read_text(encoding="utf-8", errors="replace")
        elif ext == DOCX_EXT:
            text = _read_docx(target)
        elif ext == PDF_EXT:
            text = _read_pdf(target)
        elif ext == PPTX_EXT:
            text = _read_pptx(target)
        else:
            return {"path": rel_path, "error": f"Unsupported extension {ext}"}
    except Exception as e:
        log.warning("Failed to read %s: %s", target, e)
        return {"path": rel_path, "error": str(e)}

    truncated = False
    if len(text) > max_chars:
        text = text[:max_chars]
        truncated = True
    return {
        "path": rel_path,
        "ext": ext,
        "truncated": truncated,
        "content": text,
    }

def _read_docx(path: Path) -> str:
    try:
        import docx  # python-docx
    except ImportError:
        return "[python-docx not installed; cannot read .docx]"
    doc = docx.Document(str(path))
    parts: List[str] = []
    for para in doc.paragraphs:
        parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)

def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader  # optional
    except ImportError:
        return "[pypdf not installed; cannot read .pdf]"
    reader = PdfReader(str(path))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)

def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "[python-pptx not installed; cannot read .pptx]"
    prs = Presentation(str(path))
    parts: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {idx} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)

def corpus_summary(corpus_dir: Optional[str]) -> Dict[str, Any]:
    docs = list_documents(corpus_dir)
    total = len(docs)
    by_ext: Dict[str, int] = {}
    for d in docs:
        by_ext[d["ext"]] = by_ext.get(d["ext"], 0) + 1
    return {
        "corpus_dir": corpus_dir,
        "document_count": total,
        "by_extension": by_ext,
    }
